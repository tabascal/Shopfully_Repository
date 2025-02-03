import streamlit as st
import pandas as pd
import pptx
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
import os
import time
import zipfile
import io
import shutil
from datetime import datetime
import re
import subprocess
import smtplib
import ssl
from email.message import EmailMessage
import requests

MAILCHIMP_API_KEY = "TU_API_KEY_AQUI"
MAILCHIMP_SERVER_PREFIX = "us21"  # Cambia esto según tu cuenta
MAILCHIMP_AUDIENCE_ID = "TU_LIST_ID_AQUI"

def send_email_mailchimp(receiver_email, subject, body):
    """Envía un correo usando la API de Mailchimp."""
    url = f"https://{MAILCHIMP_SERVER_PREFIX}.api.mailchimp.com/3.0/messages"
    
    headers = {
        "Authorization": f"Bearer {MAILCHIMP_API_KEY}",
        "Content-Type": "application/json"
    }
    
    data = {
        "from_email": "tu_correo@mailchimp.com",
        "subject": subject,
        "content": [{"type": "text/plain", "value": body}],
        "recipients": [{"email": receiver_email}]
    }

    response = requests.post(url, json=data, headers=headers)
    
    if response.status_code == 200:
        print(f"✅ Email sent to {receiver_email}")
    else:
        print(f"❌ Failed to send email to {receiver_email}: {response.text}")


def send_email(receiver_email, zip_file_path, subject, body, sender_email, sender_password):
    """Envía un correo con el ZIP adjunto."""
    msg = EmailMessage()
    msg["From"] = sender_email
    msg["To"] = receiver_email
    msg["Subject"] = subject
    msg.set_content(body)

    # Adjuntar el archivo ZIP
    with open(zip_file_path, "rb") as attachment:
        msg.add_attachment(
            attachment.read(),
            maintype="application",
            subtype="zip",
            filename=os.path.basename(zip_file_path),
        )

    # Configurar conexión SMTP
    context = ssl.create_default_context()
    with smtplib.SMTP_SSL("smtp.gmail.com", 465, context=context) as server:
        server.login(sender_email, sender_password)
        server.send_message(msg)

    print(f"✅ Email sent to {receiver_email}")


def convert_pptx_to_pdf(pptx_path, pdf_path):
    """Convierte un archivo PPTX a PDF en Linux usando LibreOffice (funciona en Streamlit Cloud)."""
    try:
        subprocess.run(["libreoffice", "--headless", "--convert-to", "pdf", pptx_path, "--outdir", os.path.dirname(pdf_path)], check=True)
    except Exception as e:
        print(f"Error converting {pptx_path} to PDF: {e}")


def create_zip_of_presentations(folder_path):
    """Crea un archivo ZIP con todos los PPTX generados en la carpeta."""
    zip_buffer = io.BytesIO()

    with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zipf:
        for file in os.listdir(folder_path):
            file_path = os.path.join(folder_path, file)
            if file.endswith(".pptx"):  # Evitamos incluir plantilla y Excel
                zipf.write(file_path, arcname=file)

    zip_buffer.seek(0)
    return zip_buffer


def get_filename_from_selection(row, selected_columns):
    """Genera el nombre del archivo según las columnas seleccionadas."""
    file_name_parts = [str(row[col]) for col in selected_columns if col in row]
    return "_".join(file_name_parts)


def update_text_of_textbox(presentation, column_letter, new_text):
    """Busca y reemplaza texto dentro de las cajas de texto que tengan el formato {A}, {B}, etc."""
    pattern = rf"\{{{
        column_letter}\}}"  # Expresión regular para encontrar "{A}", "{B}", etc.

    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text:
                if re.search(pattern, shape.text):  # Buscar patrón en el texto
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.text = re.sub(pattern, str(
                                new_text), run.text)  # Reemplazo



def process_files(ppt_file, excel_file, search_option, start_row, end_row, store_ids, selected_columns, output_format, use_mailchimp):
    """Genera reportes y envía emails usando Mailchimp si está activado."""
    
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    folder_name = f"Presentations_{timestamp}"
    os.makedirs(folder_name, exist_ok=True)

    ppt_template_path = os.path.join(folder_name, ppt_file.name)
    excel_file_path = os.path.join(folder_name, excel_file.name)

    with open(ppt_template_path, "wb") as f:
        f.write(ppt_file.getbuffer())
    with open(excel_file_path, "wb") as f:
        f.write(excel_file.getbuffer())

    df = pd.read_excel(excel_file_path, sheet_name=0)
    total_files = len(df)

    progress_bar = st.progress(0)
    progress_text = st.empty()

    for index, row in df.iterrows():
        email = row["Email"] if "Email" in df.columns else None
        if not email:
            st.warning(f"No email found for row {index}, skipping...")
            continue

        process_row(ppt_template_path, row, df, index, selected_columns, folder_name, output_format)
        
        if use_mailchimp:
            send_email_mailchimp(email, "Your Report is Ready", "Please find your report attached.")

    zip_path = f"{folder_name}.zip"
    shutil.make_archive(zip_path.replace(".zip", ""), 'zip', folder_name)

    with open(zip_path, "rb") as zip_file:
        st.download_button(
            label=f"📥 Download {total_files} reports ({output_format})",
            data=zip_file,
            file_name=f"{folder_name}.zip",
            mime="application/zip"
        )

    st.success("✅ All reports have been generated and emails sent!")




def process_row(presentation_path, row, df1, index, selected_columns, output_folder, output_format):
    """Procesa una fila y genera un archivo PPTX o PDF en Streamlit Cloud."""
    presentation = pptx.Presentation(presentation_path)

    for col_idx, col_name in enumerate(row.index):
        column_letter = chr(65 + col_idx)
        update_text_of_textbox(presentation, column_letter, row[col_name])

    file_name = get_filename_from_selection(row, selected_columns)
    pptx_path = os.path.join(output_folder, f"{file_name}.pptx")

    # Guardar como PPTX
    presentation.save(pptx_path)

    # Si el usuario elige PDF, convertir el archivo
    if output_format == "PDF":
        pdf_path = os.path.join(output_folder, f"{file_name}.pdf")
        convert_pptx_to_pdf(pptx_path, pdf_path)
        os.remove(pptx_path)  # Eliminar el PPTX original para solo guardar el PDF



# ========= 💡 Estilos para mejorar el diseño =========
st.markdown("""
    <style>
    div.stButton > button {
        width: 100%;
        height: 50px;
        font-size: 16px;
        border-radius: 10px;
        font-weight: bold;
    }
    </style>
""", unsafe_allow_html=True)

# ========= Título =========
st.title("Shopfully Dashboard Generator")

# Opción para elegir el formato de salida
st.markdown("### **Select Output Format**")
output_format = st.radio("Choose the file format:", ["PPTX", "PDF"])

# Mensaje de advertencia si el usuario elige PDF
if output_format == "PDF":
    st.warning("⚠️ Converting to PDF may take extra time. Large batches of presentations might take several minutes.")



# ========= 📂 Upload de archivos con formato mejorado =========
st.markdown(
    "**Upload PPTX Template**  \n*(Text Box format that will be edited -> {Column Letter} For Example: `{A}`)*", unsafe_allow_html=True)
ppt_template = st.file_uploader("", type=["pptx"])

st.write("")  # Espaciado

st.markdown(
    "**Upload Excel File**  \n*(Column A must be `Store ID`)*", unsafe_allow_html=True)
data_file = st.file_uploader("", type=["xlsx"])


# ========= 🔍 Botones mejorados para "Search by" =========
st.markdown("### **Search by:**")  # Título en negrita y más grande
col1, col2 = st.columns(2)  # Dos columnas para alinear botones en mosaico

# Inicializar la variable de estado para la selección del filtro
if "search_option" not in st.session_state:
    st.session_state.search_option = "rows"  # Valor por defecto

# Botón 1 - Search by Rows
with col1:
    if st.button("🔢 Rows", use_container_width=True):
        st.session_state.search_option = "rows"

# Botón 2 - Search by Store ID
with col2:
    if st.button("🔍 Store ID", use_container_width=True):
        st.session_state.search_option = "store_id"

# Mostrar la opción seleccionada
st.markdown(f"**Selected: `{st.session_state.search_option}`**")


# ========= 🔢 Inputs para definir el rango de búsqueda =========
start_row, end_row, store_ids = None, None, None

if st.session_state.search_option == "rows":
    start_row = st.number_input("Start Row", min_value=0, step=1)
    end_row = st.number_input("End Row", min_value=0, step=1)

elif st.session_state.search_option == "store_id":
    store_ids = st.text_input("Enter Store IDs (comma-separated)")


# ========= 📝 Selección de columnas para el nombre del archivo =========
if data_file is not None:
    # Leer la primera hoja del Excel
    df = pd.read_excel(data_file, sheet_name=0)
    column_names = df.columns.tolist()

    selected_columns = st.multiselect(
        "📂 Select and order the columns for the file name:",
        column_names,
        default=column_names[:3]
    )

    def get_filename_from_selection(row, selected_columns):
        """Genera el nombre del archivo según las columnas seleccionadas."""
        file_name_parts = [str(row[col])
                           for col in selected_columns if col in row]
        return "_".join(file_name_parts)

    st.write("🔹 Example file name:", get_filename_from_selection(
        df.iloc[0], selected_columns))


# ========= 🚀 Botón de procesamiento =========
st.markdown("### **Email Configuration**")
sender_email = st.text_input("📧 Your Email (Gmail or Outlook)")
sender_password = st.text_input("🔑 Your Email Password", type="password")
send_email_option = st.checkbox("📩 Send reports via email")

use_mailchimp = st.checkbox("📩 Send reports via Mailchimp")

if st.button("Process & Send Emails" if use_mailchimp else "Process"):
    if ppt_template and data_file:
        process_files(ppt_template, data_file, st.session_state.search_option, start_row, end_row, store_ids, selected_columns, output_format, use_mailchimp)
    else:
        st.error("Please upload files before processing.")



